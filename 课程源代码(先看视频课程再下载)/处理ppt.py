'''操作PPT使用python第三方库python-pptx'''
'''安装python-pptx'库     pip install python-pptx -i https://pypi.mirrors.ustc.edu.cn/simple/ 
   安装之后会包含依赖库：lxml        Pillow        XlsxWriter'''

'''创建一个ppt文件
   注：只是一个PowerPoint文件，其中没有任何幻灯片
   '''
# from pptx import Presentation
# prs = Presentation()   #创建演示文稿(PPT文件)对象，类似open()
# prs.save('test.pptx')  #命名为test.pptx


'''打开一个ppt文件
    注：PowerPoint 2003及更早版本中的.ppt文件将不起作用'''
# from pptx import Presentation
# prs = Presentation('test.pptx')   #创建演示文稿(PPT文件)对象，类似open()
# prs.save('test1.pptx')  #命名为test.pptx


'''创建ppt中的幻灯片  
    幻灯片的版式11个'''
# from pptx import Presentation
# prs = Presentation()   #创建演示文稿对象
# slide = prs.slides.add_slide(prs.slide_layouts[2]) #prs.slides是演示文稿中幻灯片的集合，创建幻灯片，且版式从0-10
# slide = prs.slides.add_slide(prs.slide_layouts[10]) #prs.slides是演示文稿中幻灯片的集合，创建幻灯片，且版式从0-10
# slide = prs.slides.add_slide(prs.slide_layouts[4]) #prs.slides是演示文稿中幻灯片的集合，创建幻灯片，且版式从0-10
# prs.save('test.pptx')   #创建多个幻灯片
# print(len(prs.slides))  #len可以用于测试具有ppt中幻灯片的个数


'''输入文本内容 '''
# from pptx import Presentation
# prs = Presentation()   #创建演示文稿对象
# slide = prs.slides.add_slide(prs.slide_layouts[0])
# slide.placeholders[0].text = "你好, 同学!"   #Placeholders占位符
# slide.placeholders[1].text =  "这里开始学习python-pptx!"
# prs.save('test.pptx')
# for p in slide.placeholders:
#     print('id:',p.placeholder_format.idx, '名字：',p.name)

'''占位符：可以将内容放入其中的 预 格式化 容器
用途：设置一套PPT模板，模板中提供固定格式的占位符,通过程序代码批量
      写入内容，实现批量演示文稿的开发。例如给不同客户不同PPT内容。
'''
'''PPT中除了占位符，还会存在一些其他的不能输入内容的形状，
   访问PPT中所有的形状shapes(包含占位符和其他形状)称位形状数
   slide.shapes   幻灯片中的形状树(包含占位符)
   slide.placeholders   幻灯片中的占位符'''
# from pptx import Presentation
# prs = Presentation('test.pptx')
# slide = prs.slides[0]
# for shape in slide.shapes:
#     print('第一个幻灯片：',shape.shape_type)
#
# slide_1 = prs.slides[1]
# for shape in slide_1.shapes:
#     print('第二个幻灯片：',shape.shape_type)
'''可支持的所有形状：
        EnumMember("AUTO_SHAPE", 1, "AutoShape"),
        EnumMember("CALLOUT", 2, "Callout shape"),
        EnumMember("CANVAS", 20, "Drawing canvas"),
        EnumMember("CHART", 3, "Chart, e.g. pie chart, bar chart"),
        EnumMember("COMMENT", 4, "Comment"),
        EnumMember("DIAGRAM", 21, "Diagram"),
        EnumMember("EMBEDDED_OLE_OBJECT", 7, "Embedded OLE object"),
        EnumMember("FORM_CONTROL", 8, "Form control"),
        EnumMember("FREEFORM", 5, "Freeform"),
        EnumMember("GROUP", 6, "Group shape"),
        EnumMember("IGX_GRAPHIC", 24, "SmartArt graphic"),
        EnumMember("INK", 22, "Ink"),
        EnumMember("INK_COMMENT", 23, "Ink Comment"),
        EnumMember("LINE", 9, "Line"),
        EnumMember("LINKED_OLE_OBJECT", 10, "Linked OLE object"),
        EnumMember("LINKED_PICTURE", 11, "Linked picture"),
        EnumMember("MEDIA", 16, "Media"),
        EnumMember("OLE_CONTROL_OBJECT", 12, "OLE control object"),
        EnumMember("PICTURE", 13, "Picture"),
        EnumMember("PLACEHOLDER", 14, "Placeholder"),
        EnumMember("SCRIPT_ANCHOR", 18, "Script anchor"),
        EnumMember("TABLE", 19, "Table"),
        EnumMember("TEXT_BOX", 17, "Text box"),
        EnumMember("TEXT_EFFECT", 15, "Text effect"),
        EnumMember("WEB_VIDEO", 26, "Web video"))
'''

'''使用形状树'''
# from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE
# from pptx.util import Inches   #英寸
# prs = Presentation('test.pptx')
# slide = prs.slides[0]
# left = top =  Inches(1)
# width = height = Inches(0.5)
# shape = slide.shapes.add_shape(   #通过阅读源代码
#     MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
# )
# prs.save('test.pptx')

'''add_shape的使用方法
add_shape(autoshape_type_id, left, top, width, height)
参数autoshape_type_id：自选形状的类型，在(pptx/enum/shapes)里面调
参数left：距离幻灯片左边框距离
参数top：距离幻灯片顶部边框距离
参数width：形状的宽度
参数height：形状的高度
'''

'''读取shape的位置及大小'''
# from pptx import Presentation
# from pptx.util import Cm,Inches
# prs = Presentation('test.pptx')
# slide = prs.slides[0]
# for shape in slide.shapes:
#     print(shape.shape_type)
#     print(shape.left.cm) #获取shape的到幻灯片左边距
#     print(shape.top.cm)  #默认长度单位为point，可转换cm
#     print(shape.width.cm)
#     print(shape.height.cm)


'''设置autoshape颜色（轮廓颜色和填充颜色）及文本'''
from pptx import Presentation
from pptx.dml.color import  RGBColor
from pptx.util import Pt
#
# prs = Presentation('test.pptx')
# slide = prs.slides[0]
# print(slide.shapes)  #形状树对象，其内容可以遍历、索引
# print(slide.shapes[0])  #每个形状都有自己的编码
# print(slide.shapes[2].shape_type)
# shape_3 = slide.shapes[2]   #将第3个shape对象命名为shape_3
# shape_3.text = '颜色'       #形状中添加文字
# shape_3.fill.solid()        #将填充类型设置为纯色，即纯色。注意调用此方法不设置颜色，会导致形状以纯色填充显示；相反，它启用后续指定给诸如fore_color之类的属性来设置颜色。
# shape_3.fill.fore_color.rgb = RGBColor(255, 0, 0)  #填充为红色
# shape_3.fill.background()  #透明颜色
# line = shape_3.line        #对shape的外框线进行设置
# line.color.rgb = RGBColor(255, 0, 0)  #line.fill.background()将会设置为透明色
# line.width = Pt(2.5)
# prs.save('test.pptx')



'''图片占位符'''
# from pptx import Presentation
# prs = Presentation('test.pptx')
# slide = prs.slides[1]
# print(slide.placeholders[1].name)  #不同的占位符填充不同类型的数据信息
# slide.placeholders[1].insert_picture('lufei.png')#图片填充满占位符，批量生成有利
# prs.save('test.pptx')


'''判断形状中是否可以添加文本'''
# from pptx import Presentation
# prs = Presentation('test.pptx')
# slide = prs.slides[1]
# for shape in slide.shapes:
#     if shape.has_text_frame:   #判断形状中是否包含文本框架，如果包含，则返回True
        # print(shape.name,'可以输入文本',shape.shape_id)
        # shape.text = '内容'
        # shape.text_frame.clear()   #清除文本框架中的内容
# prs.save('test.pptx')


'''包含文本框架的形状至少包含一个段落对象
    段落对象设置方法与word文档中段落方法相同'''
# from pptx import Presentation
# from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
#
# prs = Presentation('test.pptx')
# slide = prs.slides[1]
# biaozhu = slide.shapes[5]           #将标注shape对象赋值给变量biaozhu
# print(biaozhu.text_frame.paragraphs[0]) #获取标注shape的段落
# print(biaozhu.text_frame.paragraphs[0].text) #获取标注shape的段落

# new_run = biaozhu.text_frame.paragraphs[0].add_run()
# new_run.text  = '附加一个新内容'
# new_par = biaozhu.text_frame.add_paragraph()  #在标注shape中添加一个段落
# new_par.text = '我是新添加的段落'   #设置新添加的段落内容
# new_par.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
# prs.save('test.pptx')    #   ctrl +alt +⬅  表示跳回上一步  ctrl+鼠标点击


'''绘制图表'''
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

prs = Presentation('test.pptx')
blank = prs.slide_layouts[6]    #设置一个幻灯片版式
slide = prs.slides.add_slide(blank) #添加幻灯片

chart_data = CategoryChartData()  #创建一个图表对象
chart_data.categories = ['美国', '英国', '印度','法国']   #图表对象中添加x轴的内容
chart_data.add_series('国外疫情', (8874793, 674644,7809899,823786)) #添加数据
x, y, cx, cy = Inches(2), Inches(1), Inches(5), Inches(6)
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED,x,y,cx,cy,chart_data)
prs.save('test.pptx')


'''add_series的使用
 add_series(name, values=(), number_format=None)
参数name：图表的名称，在幻灯片中名称是唯一的
参数values：图表的数据信息，与categories中的一一对应
参数number_format：指定数据的显示方式
'''

'''add_chart的使用
add_chart(chart_type, x, y, cx, cy, chart_data)
参数chart_type：显示的图表类型，在pptx/enum/chart中，例如饼状图、条形图等等
参数x，y分别表示图表的位置
参数cx，cy分别表示图表的大小
参数chart_data：插入图表的数据
'''

'''插入表格'''
# from pptx import Presentation
# from pptx.util import Inches
# prs = Presentation('test.pptx')
# slide = prs.slides.add_slide(prs.slide_layouts[5])
# x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
# shape = slide.shapes.add_table(3, 3, x, y, cx, cy)   #添加一个表格框架
# prs.save('test.pptx')
#注：add_table 和add_chart都是在形状数shapes里面进行添加一个形状
#因此add_table是包含表格的形状
'''add_table的使用方法
add_table(rows, cols, left, top, width, height)
参数rows：表格的行数
参数cols：表格的列数
参数left，top：指定表格的位置
参数width，height指定表格的宽度和高度，并且会均匀的分布在每一个行和列中
'''

'''范围表格cell'''
# from pptx import Presentation
# prs = Presentation('test.pptx')
# slide = prs.slides[3]
# print(slide.shapes[1].shape_type)
# table = slide.shapes[1].table
# print(table.cell(0,0).text)   #excel cell
# table.cell(1,0).text = '张三'   #进入cell中，可以查阅相关操作，例如文本设置，具备文本框架
# prs.save('test.pptx')


'''项目实例：实现公司营收统计创建ppt'''
import openpyxl
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches,Pt
ppt_dir = './ppt项目--公司营收统计/'

def slide_0_fit(s_0,file_name):
    '''功能幻灯片0的设置'''
    s_01 = s_0.shapes[1]
    ph = s_01.text_frame.paragraphs[0]
    ph.text = file_name[0:-5]
    s_01.text_frame.fit_text(max_size=88, bold=True)

    s_02 = s_0.shapes[2]
    ph = s_02.text_frame.paragraphs[0]
    ph.text = '北京+上海+深圳+四川'
    s_02.text_frame.fit_text(max_size=40)

def get_data(sheet_file,start,end):
    '''获取对应excel中位置的数据为列表'''
    data = []
    for line in sheet_file[start:end]:
        for cell in line:
            data.append(cell.value)
    return data

def slide_3_fit(s_3,sheet_file,sheet_name):
    '''功能：幻灯片3，4，5，6的设置'''
    s_3.shapes[0].text = sheet_name

    chart_data = CategoryChartData()  #创建一个图表对象
    riqi = get_data(sheet_file,'A2','A31')
    xinpian = get_data(sheet_file,'B2','B31')
    shouji = get_data(sheet_file,'C2','C31')
    zhineng = get_data(sheet_file,'D2','D31')
    chart_data.categories = riqi   #图表对象中添加x轴的内容
    chart_data.add_series('芯片', tuple(xinpian)) #添加数据
    chart_data.add_series('手机', tuple(shouji)) #添加数据
    chart_data.add_series('智能设备', tuple(zhineng)) #添加数据
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(10), Inches(5)
    s_3.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED,x,y,cx,cy,chart_data)

def ppt_create(ex_file,file_name):
    '''PPT创建'''
    prs = Presentation('./ppt项目--公司营收统计/模板.pptx')
    slide_0_fit(prs.slides[0],file_name)  #设置幻灯片1

    sheet_names = ex_file.get_sheet_names()         #获取工作簿中的所有工作表名
    for sheet_name in sheet_names:
        sheet_file = ex_file.get_sheet_by_name(sheet_name)  # 读取到对应工作表
        '''幻灯片3，4，5的处理'''
        slide_3 = prs.slides.add_slide(prs.slide_layouts[5])
        slide_3_fit(slide_3,sheet_file,sheet_name)
    prs.save(ppt_dir+file_name[:-5]+'.pptx')

'''项目开始'''
for file_name in os.listdir('./excel项目-##公司15年营收'):
    print(file_name)
    ex_file_name = './excel项目-##公司15年营收/'+file_name
    print(ex_file_name)
    ex_file = openpyxl.load_workbook(ex_file_name)   #打开excel_file_name文件
    ppt_create(ex_file,file_name)

